import streamlit as st
import google.generativeai as genai
from groq import Groq
import uuid, time, os
from dotenv import load_dotenv
from pypdf import PdfReader
from pptx import Presentation

# --- HEAVY IMPORTS ---
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from ddgs import DDGS # 💡 Naya aur Fast Web Search Tool

# Load API Keys
load_dotenv()

# --- 1. SESSION STATE INITIALIZATION ---
if "score" not in st.session_state: st.session_state.score = 0
if "total" not in st.session_state: st.session_state.total = 0
if "chats" not in st.session_state: st.session_state.chats = {}
if "active_chat" not in st.session_state: st.session_state.active_chat = None
if "retriever" not in st.session_state: st.session_state.retriever = None

# --- 2. CORE CLIENTS ---
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))

# --- 3. UI STYLING (Premium Dark Theme) ---
st.set_page_config(page_title="Agentic Quiz Master Pro", layout="wide", page_icon="🎓")

st.markdown("""
    <style>
        .stApp { background: radial-gradient(circle at top left, #1e2229, #131314); color: #e3e3e3; }
        [data-testid="stSidebar"] { background: rgba(30, 31, 32, 0.7) !important; backdrop-filter: blur(12px); border-right: 1px solid rgba(255, 255, 255, 0.1); }
        .stButton>button { width: 100%; border-radius: 12px; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; border: none; font-weight: 600; margin-top: 10px; }
        [data-testid="stChatMessage"] { background: rgba(255, 255, 255, 0.04); border: 1px solid rgba(255, 255, 255, 0.08); border-radius: 15px !important; }
        [data-testid="stMetric"] { background: rgba(255, 255, 255, 0.05); padding: 10px; border-radius: 10px; border: 1px solid rgba(255, 255, 255, 0.1); }
    </style>
""", unsafe_allow_html=True)

# --- 4. LOGIC: FILE PROCESSING (Local Knowledge) ---
def process_file(uploaded_file):
    try:
        text = ""
        if uploaded_file.name.endswith('.pdf'):
            reader = PdfReader(uploaded_file)
            text = "".join([p.extract_text() for p in reader.pages if p.extract_text()])
        elif uploaded_file.name.endswith('.pptx'):
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + " "

        if not text.strip():
            st.error("❌ File khali hai ya readable text nahi mila!")
            return None

        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=150)
        chunks = text_splitter.split_text(text)
        
        embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
        vectorstore = FAISS.from_texts(chunks, embeddings)
        return vectorstore.as_retriever()
    except Exception as e:
        st.error(f"⚠️ Error: {e}")
        return None

# --- 5. ORCHESTRATION LAYER (THE AI ROUTER & WEB SEARCH) ---
def route_query(prompt):
    """Faisla karta hai ke PDF mein dekhna hai ya Internet par."""
    router_prompt = f"""
    You are an intelligent routing agent. Classify the user's query.
    If the user asks about "latest news", "today", "yesterday", "current events", or general knowledge outside a static PDF, output exactly: WEB
    Otherwise, if it seems related to studying an uploaded document or a general quiz, output exactly: PDF
    
    User Query: {prompt}
    """
    try:
        # Router ab Groq par chal raha hai (Fast & Reliable)
        res = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": router_prompt}]
        )
        decision = res.choices[0].message.content.strip().upper()
        return "WEB" if "WEB" in decision else "PDF"
    except Exception as e:
        return "PDF" 

def get_web_context(prompt):
    """Internet se real-time data fetch karta hai using strict smart keywords."""
    try:
        # AI se strict instructions de kar keywords nikalwana
        query_res = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": f"You are a search query generator. Extract 2 to 4 keywords from the user prompt for a Google search. Output ONLY the keywords. NO quotes, NO extra words, NO introduction. Prompt: {prompt}"}]
        )
        
        # Cleaner: Extra quotation marks aur spaces hatana
        search_keywords = query_res.choices[0].message.content.replace('"', '').replace("'", "").strip()
        
        # 💡 TERMINAL TRACKER: Yeh aapko VS Code mein dikhayega ke search kya ho raha hai
        print(f"🤖 AGENT SEARCHING FOR: {search_keywords}")
        
        results = DDGS().text(search_keywords, max_results=5)
        context = ""
        
        if results:
            for r in results:
                context += r.get('body', '') + "\n"
                
        # FALLBACK (Plan B): Agar pehli search fail ho jaye toh default news laye
        if not context.strip():
            print("⚠️ Smart search failed, trying backup search...")
            backup_results = DDGS().text("latest global technology news", max_results=3)
            for r in backup_results:
                context += r.get('body', '') + "\n"
                
        return context if context.strip() else "No information found on the internet."
    except Exception as e:
        return f"Web Search Failed: {e}"

# --- 6. AI ENGINE (Upgraded Agentic Output) ---
QUIZ_MASTER_PROMPT = """
You are a STRICT and expert Exam Master. 
- Generate MCQs based strictly on the provided context.
- If evaluating an answer, start response with 'CORRECT!' if user is right, or 'INCORRECT.' if wrong. Do NOT be polite if they are wrong.
- Provide a brief explanation.
- IMPORTANT: If the context is from the WEB, mention that it is based on live web data.
- ALWAYS check user answers against the PREVIOUS CHAT MEMORY.
"""

def get_quiz_response(prompt, chat_messages, retriever=None):
    context = ""
    source_label = ""
    
    # 💡 The Agentic Decision (Where to look?)
    route = route_query(prompt)
    
    if route == "WEB":
        context = get_web_context(prompt)
        source_label = "Live Web Search 🌐"
    else:
        if retriever:
            try:
                docs = retriever.invoke(prompt)
                context = "\n".join([d.page_content for d in docs])
                source_label = "Uploaded PDF 📄"
            except: 
                context = "No document uploaded."
                source_label = "None"
        else:
            context = "No document uploaded."
            source_label = "None"

    # Memory Building
    history_text = ""
    for msg in chat_messages[-5:-1]: 
        history_text += f"{msg['role'].upper()}: {msg['content']}\n"

    full_p = f"{QUIZ_MASTER_PROMPT}\n\nPREVIOUS CHAT MEMORY:\n{history_text}\n\nCONTEXT FROM {source_label}:\n{context}\n\nUSER: {prompt}"
    
    try:
        model = genai.GenerativeModel('gemini-1.5-flash')
        return model.generate_content(full_p).text, "Gemini Engine", source_label
    except:
        res = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "system", "content": QUIZ_MASTER_PROMPT}, {"role": "user", "content": full_p}]
        )
        return res.choices[0].message.content, "Groq Backup", source_label

# --- 7. SIDEBAR UI ---
with st.sidebar:
    st.markdown("<h1 style='color: #8ab4f8; text-align: center;'>🤖 Agentic Quiz Master</h1>", unsafe_allow_html=True)
    
    st.markdown("### 📊 Live Stats")
    c1, c2 = st.columns(2)
    c1.metric("Correct ✅", st.session_state.score)
    c2.metric("Total 📝", st.session_state.total)
    
    if st.button("Reset Stats"):
        st.session_state.score = 0
        st.session_state.total = 0
        st.rerun()

    st.write("---")
    file = st.file_uploader("Upload PDF or PPTX", type=["pdf", "pptx"])
    if file and st.button("🚀 Process Material"):
        with st.spinner("AI is reading and indexing..."):
            st.session_state.retriever = process_file(file)
            if st.session_state.retriever: st.success("Knowledge Loaded!")

    if st.button("➕ Start New Session", use_container_width=True):
        cid = str(uuid.uuid4())[:8]
        st.session_state.chats[cid] = {"messages": []}
        st.session_state.active_chat = cid
        st.rerun()

# --- 8. MAIN CHAT INTERFACE ---
if st.session_state.active_chat:
    active_cid = st.session_state.active_chat
    chat_history = st.session_state.chats[active_cid]
    
    for m in chat_history["messages"]:
        with st.chat_message(m["role"]): st.markdown(m["content"])

    if prompt := st.chat_input("Ask a question, request a PDF quiz, or ask about today's news..."):
        chat_history["messages"].append({"role": "user", "content": prompt})
        with st.chat_message("user", avatar="👤"): st.markdown(prompt)

        with st.chat_message("assistant", avatar="🤖"):
            with st.spinner("Agent is thinking & searching..."):
                res, engine, source = get_quiz_response(prompt, chat_history["messages"], st.session_state.retriever)
            
            # SCORING LOGIC
            if res.strip().startswith("CORRECT!"):
                st.session_state.score += 1
                st.session_state.total += 1
                st.balloons()
            elif res.strip().startswith("INCORRECT."):
                st.session_state.total += 1

            st.markdown(res)
            # Displaying both Engine and Source
            st.caption(f"⚙️ Engine: {engine} | 📚 Source: {source}")
            chat_history["messages"].append({"role": "assistant", "content": res})
            
else:
    st.markdown("<div style='text-align: center; margin-top: 100px;'>", unsafe_allow_html=True)
    st.title("🤖 Welcome to Agentic Quiz Master Pro")
    st.write("Upload a file from the sidebar OR just ask a question to trigger the Web Search Agent!")
    st.markdown("</div>", unsafe_allow_html=True)